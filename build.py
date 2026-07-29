#!/usr/bin/env python3
"""
KLARTEXT — Build.

Liest folien.md, rendert jede Folie als HTML, schiesst sie mit headless Chrome
als PNG und packt alles in eine PPTX mit Referentennotizen sowie ein PDF.

    python3 build.py            # alles bauen
    python3 build.py --schnell  # nur HTML, ohne Chrome (fuer Layout-Checks)
    python3 build.py --nur 3,7  # nur Folie 3 und 7 neu schiessen

Format von folien.md siehe README-FOLIEN.md.
"""

import html
import re
import shutil
import subprocess
import sys
import time
from pathlib import Path

WURZEL = Path(__file__).parent
QUELLE = WURZEL / "folien.md"
AUS = WURZEL / "ausgabe"
STIL = WURZEL / "stil.css"
PORT = 8791

CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"

# Sprechtempo fuer die Zeitschaetzung: Woerter pro Minute, ruhiger Vortrag
WPM = 125
# Gesamtfenster 22 Minuten: 18 fuer die Folien, 4 fuer die Vorfuehrung.
TOOL_MIN = 4
ZIEL_SEK = 18 * 60


# --------------------------------------------------------------------------
# Parser
# --------------------------------------------------------------------------

def entkleide(wert):
    """Umschliessende Anfuehrungszeichen entfernen. Sie sind optional und
    dienen nur dazu, Werte mit Doppelpunkt eindeutig zu machen."""
    w = wert.strip()
    if len(w) >= 2 and w[0] == w[-1] and w[0] in ('"', "'"):
        return w[1:-1].strip()
    return w


def parse(text):
    """Zerlegt folien.md in eine Liste von Folien-Dicts."""
    folien = []
    aktuell = None
    modus = None          # None | "kopf" | "notiz"
    liste_key = None

    for nr, zeile in enumerate(text.splitlines(), 1):
        roh = zeile.rstrip()

        # Neue Folie
        if roh.startswith("## "):
            if aktuell:
                folien.append(aktuell)
            aktuell = {"_titel_kommentar": roh[3:].strip(), "_zeile": nr, "notiz": []}
            modus, liste_key = "kopf", None
            continue

        if aktuell is None:
            continue      # Vorspann vor der ersten Folie wird ignoriert

        # Notizblock
        if roh.strip().upper() in ("### NOTIZ", "### NOTIZEN"):
            modus, liste_key = "notiz", None
            continue

        if modus == "notiz":
            aktuell["notiz"].append(roh)
            continue

        # Kopfbereich
        if not roh.strip():
            liste_key = None
            continue

        if roh.lstrip().startswith("- ") and liste_key:
            aktuell[liste_key].append(entkleide(roh.lstrip()[2:]))
            continue

        m = re.match(r"^([a-zA-ZäöüÄÖÜ_][\w\-äöüÄÖÜß]*)\s*:\s*(.*)$", roh)
        if m:
            key, wert = m.group(1).strip(), entkleide(m.group(2))
            if wert == "":
                aktuell[key] = []
                liste_key = key
            else:
                aktuell[key] = wert
                liste_key = None
        else:
            raise SyntaxError(
                f"folien.md Zeile {nr}: verstehe ich nicht.\n  >>> {roh}\n"
                f"  Erwartet: 'schluessel: wert', '- Listenpunkt', '## Neue Folie' "
                f"oder '### NOTIZ'."
            )

    if aktuell:
        folien.append(aktuell)

    for f in folien:
        f["notiz"] = "\n".join(f["notiz"]).strip()
    return folien


# --------------------------------------------------------------------------
# Hilfen fuer die Ausgabe
# --------------------------------------------------------------------------

def e(s):
    """HTML-escapen, aber **fett** und *kursiv* durchlassen."""
    s = html.escape(str(s or ""))
    s = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", s)
    s = re.sub(r"(?<!\*)\*([^*]+?)\*(?!\*)", r"<i>\1</i>", s)
    return s


def kopf(f, seite, gesamt):
    # Backup-Folien zaehlen nicht zum Hauptteil und tragen keine Seitenzahl
    pg = "BACKUP" if f.get("backup") else f"{seite:02d} / {gesamt:02d}"
    return f"""<div class="head">
  <div class="mark"><span class="sq"></span><span class="w">KLARTEXT</span></div>
  <div class="sec">{e(f.get('kapitel',''))}</div>
  <div class="pg">{pg}</div>
</div>
<div class="headrule"></div>"""


def fuss(f):
    links = f.get("fussl", "vhs Frankfurt · Projekt KLARTEXT")
    rechts = f.get("fussr", "Henrik Heil · cimdata 2026")
    return f'<div class="foot"><span>{e(links)}</span><span>{e(rechts)}</span></div>'


def h1(f):
    t = e(f.get("titel", ""))
    if f.get("akzent"):
        t += f' <span class="g">{e(f["akzent"])}</span>'
    klasse = " klein" if f.get("klein") else ""
    out = f'<h1 class="{klasse.strip()}">{t}</h1>'
    if f.get("lede"):
        out += f'<p class="lede">{e(f["lede"])}</p>'
    return out


def callout(f):
    if not f.get("callout"):
        return ""
    kl = f'<span class="kl">{e(f["calloutsub"])}</span>' if f.get("calloutsub") else ""
    return f'<div class="callout">{e(f["callout"])}{kl}</div>'


def quellen(f):
    """Quellenzeile. Ein Eintrag der Form 'Text | URL' wird zum Link.

    Als Liste geschrieben ist jeder Listenpunkt ein Eintrag; als einzelne
    Zeile trennt ' · ' die Eintraege. Ohne URL bleibt der Eintrag Text.
    """
    if not f.get("quellen"):
        return ""
    q = f["quellen"]
    teile = q if isinstance(q, list) else str(q).split(" · ")
    stuecke = []
    for t in teile:
        t = str(t).strip()
        if not t:
            continue
        if "|" in t:
            text, url = (x.strip() for x in t.split("|", 1))
            stuecke.append(f'<a href="{e(url)}">{e(text)}</a>')
        else:
            stuecke.append(f"<span>{e(t)}</span>")
    if not stuecke:
        return ""
    # Trenner als eigenes Element, nicht als ::after im Link: sonst zieht
    # die Unterlinie des Links ueber den Punkt bis zum naechsten Eintrag
    inhalt = '<span class="qsep">·</span>'.join(stuecke)
    return ('<div class="quellen"><span class="qlabel">Quellen</span>'
            + inhalt + "</div>")


# --------------------------------------------------------------------------
# Folientypen
# --------------------------------------------------------------------------

def bildspalte(f):
    """Rechte Spalte mit Bild, nur wenn das Feld 'bild' gesetzt ist."""
    if not f.get("bild"):
        return ""
    pfad = WURZEL / f["bild"]
    if not pfad.exists():
        print(f"  ! Bild fehlt: {f['bild']}")
        return ""
    bu = f'<div class="bu">{e(f["bu"])}</div>' if f.get("bu") else ""
    return (f'<div class="bildspalte"><div class="rahmen">'
            f'<img src="../{f["bild"]}" alt=""></div>{bu}</div>')


def seite(f, s, g, inhalt):
    """Gemeinsamer Rahmen fuer alle Textfolien, mit oder ohne Bildspalte."""
    bild = bildspalte(f)
    if bild:
        body = f'<div class="body"><div class="inhalt">{inhalt}</div>{bild}</div>'
        kl = " mitbild"
    else:
        body = f'<div class="body">{inhalt}</div>'
        kl = ""
    return f"""<div class="stage{kl}">
  {kopf(f, s, g)}
  {h1(f)}
  {body}
  {quellen(f)}
  {fuss(f)}
</div>"""


def t_panel(f, seite_, gesamt, schluss=False):
    sub = f'<p class="sub">{e(f["untertitel"])}</p>' if f.get("untertitel") else ""
    # meta: mit | getrennte Angaben, gesetzt als Versalienzeile mit Trennstrichen
    meta = ""
    if f.get("meta"):
        teile = [t.strip() for t in str(f["meta"]).split("|") if t.strip()]
        meta = ('<div class="meta">'
                + "".join(f"<span>{e(t)}</span>" for t in teile)
                + "</div>")
    titel = e(f.get("titel", ""))
    if f.get("akzent"):
        titel += f' <span class="g">{e(f["akzent"])}</span>'

    hero, kl = "", ""
    if f.get("bild") and (WURZEL / f["bild"]).exists():
        hero = f'<div class="hero"><img src="../{f["bild"]}" alt=""></div>'
        kl = " mitbild"
    elif f.get("bild"):
        print(f"  ! Bild fehlt: {f['bild']}")

    return f"""<div class="panel{kl}">
  {hero}
  <div class="inhalt">
    <div class="mark"><span class="sq"></span><span class="w">KLARTEXT</span></div>
    <div class="spacer"></div>
    <div class="rule"></div>
    <h1>{titel}</h1>
    {sub}
    {meta}
    <div class="spacer"></div>
    <div class="pfoot">
      <span>{e(f.get('fussl','ABSCHLUSSPROJEKT · CHANGE UND KI'))}</span>
      <span>{e(f.get('fussr','CIMDATA · HENRIK HEIL · 2026'))}</span>
    </div>
  </div>
</div>"""


def t_kapitel(f, seite_, gesamt):
    return f"""<div class="stage">
  {kopf(f, seite_, gesamt)}
  <div class="kap">
    <div class="kapnum">{e(f.get('nummer','01'))}</div>
    <div class="kaptxt">{h1(f)}</div>
  </div>
  {fuss(f)}
</div>"""


def t_punkte(f, seite_, gesamt):
    items = []
    for p in f.get("punkte", []):
        if "||" in p:
            haupt, sub = p.split("||", 1)
            inner = f'{e(haupt.strip())}<span class="psub">{e(sub.strip())}</span>'
        else:
            inner = e(p)
        items.append(f'<div class="pkt"><div class="bar"></div>'
                     f'<div class="ptxt">{inner}</div></div>')
    eng = " eng" if len(items) >= 6 else ""
    return seite(f, seite_, gesamt,
                 f'<div class="punkte{eng}">{"".join(items)}</div>{callout(f)}')


def t_zahlen(f, seite_, gesamt):
    zs = []
    for z in f.get("zahlen", []):
        teile = [x.strip() for x in z.split("||")]
        wert = teile[0]
        lab = teile[1] if len(teile) > 1 else ""
        klassen = []
        if len(teile) > 2 and teile[2] == "warn":
            klassen.append("warn")
        # lange Werte wie "154 T€" wuerden sonst umbrechen
        if len(wert) > 8:
            klassen.append("sehrlang")
        elif len(wert) > 5:
            klassen.append("lang")
        kl = (" " + " ".join(klassen)) if klassen else ""
        zs.append(f'<div class="z"><div class="znum{kl}">{e(wert)}</div>'
                  f'<div class="zlab">{e(lab)}</div></div>')
    return seite(f, seite_, gesamt,
                 f'<div class="zahlen">{"".join(zs)}</div>{callout(f)}')


def t_zweispalt(f, seite_, gesamt):
    sp = []
    for i in (1, 2):
        kopfz = f.get(f"spalte{i}", "")
        punkte = f.get(f"punkte{i}", [])
        alt = " alt" if i == 2 else ""
        lis = "".join(f"<li>{e(p)}</li>" for p in punkte)
        sp.append(f'<div class="sp"><span class="spkopf{alt}">{e(kopfz)}</span>'
                  f'<ul>{lis}</ul></div>')
    return seite(f, seite_, gesamt,
                 f'<div class="spalten">{"".join(sp)}</div>{callout(f)}')


def t_tabelle(f, seite_, gesamt):
    kopfz = [x.strip() for x in f.get("spalten", "").split("|")]
    ths = "".join(f"<th>{e(k)}</th>" for k in kopfz)
    trs = []
    for zeile in f.get("zeilen", []):
        tds = []
        for z in [x.strip() for x in zeile.split("|")]:
            kl = ""
            # Das fuehrende + bzw. ! bleibt stehen. Faerbt man die Zelle nur
            # ein, haengt die Unterscheidung allein an der Farbe — WCAG 1.4.1,
            # Stufe A und damit verbindlich. Ausgerechnet diese Arbeit sollte
            # das nicht tun.
            if z.startswith("+"):
                kl = ' class="ja"'
            elif z.startswith("!"):
                kl = ' class="nein"'
            elif re.match(r"^[\d.,]+\s*%?$", z):
                kl = ' class="num"'
            tds.append(f"<td{kl}>{e(z)}</td>")
        trs.append(f"<tr>{''.join(tds)}</tr>")
    tab = f'<table><thead><tr>{ths}</tr></thead><tbody>{"".join(trs)}</tbody></table>'
    return seite(f, seite_, gesamt, tab + callout(f))


def t_zitat(f, seite_, gesamt):
    quelle = f'<div class="quelle">{e(f["quelle"])}</div>' if f.get("quelle") else ""
    z = f'<div class="zitat"><div class="q">{e(f.get("zitat",""))}</div>{quelle}</div>'
    return seite(f, seite_, gesamt, z + callout(f))


def t_text(f, seite_, gesamt):
    abs_ = "".join(f'<p class="lede" style="font-size:18px;color:var(--ink);'
                   f'max-width:960px">{e(a)}</p>' for a in f.get("absaetze", []))
    return seite(f, seite_, gesamt, abs_ + callout(f))



def t_matrix(f, seite_, gesamt):
    """Stakeholder-Matrix Einfluss x Betroffenheit, Quadrantennamen laut Aufgabenblatt."""
    def q(key, stark=False):
        roh = f.get(key, "")
        teile = [x.strip() for x in roh.split("||")] if roh else [""]
        titel = teile[0]
        lis = "".join(f"<li>{e(x)}</li>" for x in teile[1:])
        kl = " stark" if stark else ""
        return f'<div class="quad{kl}"><div class="qt">{e(titel)}</div><ul>{lis}</ul></div>'

    matrix = (q("oben_links") + q("oben_rechts", stark=True)
              + q("unten_links") + q("unten_rechts"))
    inhalt = f"""<div class="matrixwrap">
  <div class="yachse"><span>{e(f.get('yhoch','Einfluss hoch'))}</span>
                      <span>{e(f.get('yniedrig','niedrig'))}</span></div>
  <div class="matrixcol">
    <div class="matrix">{matrix}</div>
    <div class="xachse"><span>{e(f.get('xniedrig','Betroffenheit niedrig'))}</span>
                        <span>{e(f.get('xhoch','hoch'))}</span></div>
  </div>
</div>{callout(f)}"""
    return seite(f, seite_, gesamt, inhalt)


def t_timeline(f, seite_, gesamt):
    """Integrierte Timeline: drei Straenge ueber drei Monate."""
    kopf_ = "".join(f"<span>{e(x.strip())}</span>"
                    for x in f.get("monate", "Monat 1|Monat 2|Monat 3").split("|"))
    zeilen = []
    for i in (1, 2, 3):
        roh = f.get(f"strang{i}")
        if not roh:
            continue
        teile = [x.strip() for x in roh.split("||")]
        name, felder = teile[0], teile[1:]
        fs = []
        for j, feld in enumerate(felder):
            kl = ""
            if feld.startswith("+"):
                kl, feld = " aktiv", feld[1:].strip()
            elif feld.startswith("~"):
                kl, feld = " mint", feld[1:].strip()
            fs.append(f'<div class="tlfeld{kl}">{e(feld)}</div>')
        zeilen.append(f'<div class="tlzeile"><div class="tlname">{e(name)}</div>'
                      f'{"".join(fs)}</div>')
    inhalt = (f'<div class="tl"><div class="tlkopf">{kopf_}</div>'
              f'{"".join(zeilen)}</div>{callout(f)}')
    return seite(f, seite_, gesamt, inhalt)


TYPEN = {
    "titel":     lambda f, s, g: t_panel(f, s, g),
    "schluss":   lambda f, s, g: t_panel(f, s, g, True),
    "kapitel":   t_kapitel,
    "punkte":    t_punkte,
    "zahlen":    t_zahlen,
    "zweispalt": t_zweispalt,
    "tabelle":   t_tabelle,
    "zitat":     t_zitat,
    "text":      t_text,
    "matrix":    t_matrix,
    "timeline":  t_timeline,
}


def rendere(f, seite, gesamt):
    typ = f.get("typ", "punkte")
    if typ not in TYPEN:
        raise ValueError(f"Folie {seite}: Typ '{typ}' kenne ich nicht. "
                         f"Moeglich: {', '.join(sorted(TYPEN))}")
    inner = TYPEN[typ](f, seite, gesamt)
    return (f'<!doctype html><html lang="de"><head><meta charset="utf-8">'
            f'<link rel="stylesheet" href="../stil.css"></head>'
            f'<body>{inner}</body></html>')


# --------------------------------------------------------------------------
# Ausgabe
# --------------------------------------------------------------------------

def schiesse_pngs(folien, nur=None):
    server = subprocess.Popen(
        [sys.executable, "-m", "http.server", str(PORT)],
        cwd=WURZEL, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    time.sleep(1.0)
    try:
        for i, _ in enumerate(folien, 1):
            if nur and i not in nur:
                continue
            ziel = AUS / f"folie-{i:02d}.png"
            subprocess.run([
                CHROME, "--headless=new", "--disable-gpu", "--hide-scrollbars",
                "--force-device-scale-factor=2", "--window-size=1280,720",
                f"--screenshot={ziel}",
                f"http://localhost:{PORT}/ausgabe/folie-{i:02d}.html",
            ], check=True, capture_output=True)
            print(f"  Folie {i:02d} gerendert")
    finally:
        server.terminate()


def baue_pptx(folien):
    from pptx import Presentation
    from pptx.util import Emu
    B, H = Emu(12192000), Emu(6858000)
    prs = Presentation()
    prs.slide_width, prs.slide_height = B, H
    leer = prs.slide_layouts[6]
    for i, f in enumerate(folien, 1):
        png = AUS / f"folie-{i:02d}.png"
        if not png.exists():
            print(f"  ! Folie {i:02d}: PNG fehlt, uebersprungen")
            continue
        s = prs.slides.add_slide(leer)
        s.shapes.add_picture(str(png), 0, 0, width=B, height=H)
        s.notes_slide.notes_text_frame.text = f.get("notiz", "") or "(keine Notiz)"
    ziel = AUS / "Praesentation-KLARTEXT.pptx"
    prs.save(str(ziel))
    return ziel


def baue_pdf_links(folien):
    """PDF direkt aus dem HTML drucken statt aus den PNGs.

    Damit bleiben Text und Links erhalten: die Quellenzeilen sind anklickbar,
    der Inhalt ist durchsuchbar und fuer Vorleseprogramme lesbar. Das
    Bild-PDF aus baue_pdf() kann beides nicht.
    """
    teile = []
    for i, _ in enumerate(folien, 1):
        p = AUS / f"folie-{i:02d}.html"
        if not p.exists():
            continue
        h = p.read_text(encoding="utf-8")
        inner = h.split("<body>", 1)[1].rsplit("</body>", 1)[0]
        teile.append(f'<div class="blatt">{inner}</div>')
    if not teile:
        return None

    sammel = AUS / "alle.html"
    sammel.write_text(
        '<!doctype html><html lang="de"><head><meta charset="utf-8">'
        '<title>KLARTEXT</title>'
        '<link rel="stylesheet" href="../stil.css">'
        "<style>"
        # Chrome nimmt in @page size keine px: 1280x720 bei 96 dpi sind
        # 13.3333 x 7.5 Zoll. Mit px bleibt es bei einer einzigen A4-Seite.
        "@page{ size:13.3333in 7.5in; margin:0; }"
        # stil.css beschneidet html,body auf 720px Hoehe. Fuer die Sammelseite
        # muss das aufgehoben werden, sonst endet der Druck nach Folie 1.
        "html,body{ margin:0; padding:0; background:#fff;"
        " width:auto; height:auto; overflow:visible; }"
        ".blatt{ width:1280px; height:720px; overflow:hidden;"
        " page-break-after:always; break-after:page; }"
        ".blatt:last-child{ page-break-after:auto; break-after:auto; }"
        "</style></head><body>" + "".join(teile) + "</body></html>",
        encoding="utf-8")

    ziel = AUS / "Praesentation-KLARTEXT.pdf"
    server = subprocess.Popen(
        [sys.executable, "-m", "http.server", str(PORT)],
        cwd=WURZEL, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    time.sleep(1.0)
    try:
        subprocess.run([
            CHROME, "--headless=new", "--disable-gpu", "--no-pdf-header-footer",
            # ohne Zeitbudget druckt Chrome, bevor die Schriften geladen sind:
            # der Text bleibt dann unsichtbar und das PDF enthaelt nur Flaechen
            "--virtual-time-budget=20000",
            f"--print-to-pdf={ziel}",
            f"http://localhost:{PORT}/ausgabe/alle.html",
        ], check=True, capture_output=True)
    finally:
        server.terminate()
    return ziel


def zeit(folien):
    print("\n  Sprechzeit")
    ges = 0
    for i, f in enumerate(folien, 1):
        w = len((f.get("notiz") or "").split())
        sek = round(w / WPM * 60)
        if not f.get("backup"):
            ges += sek
        marke = "  " if sek <= 110 else " !"
        print(f"   {marke} Folie {i:02d}  {w:>4} W.  {sek//60}:{sek%60:02d}"
              f"   {f.get('_titel_kommentar','')[:40]}")
    print(f"      GESAMT      {ges//60}:{ges%60:02d} Minuten")
    # Vortragsfenster 22 Minuten, davon 4 fuer die Vorfuehrung des Werkzeugs.
    if ges > ZIEL_SEK:
        fehl = ges - ZIEL_SEK
        print(f"      ! {fehl//60}:{fehl%60:02d} ueber dem Ziel von "
              f"{ZIEL_SEK//60} Minuten Folienzeit")
    else:
        rest = ZIEL_SEK - ges
        print(f"      Puffer {rest//60}:{rest%60:02d} bis zum Ziel von "
              f"{ZIEL_SEK//60} Minuten, danach {TOOL_MIN} Minuten Werkzeug")


def commit(folien):
    if not (WURZEL / ".git").exists():
        return
    subprocess.run(["git", "add", "-A"], cwd=WURZEL, capture_output=True)
    st = subprocess.run(["git", "diff", "--cached", "--name-only"],
                        cwd=WURZEL, capture_output=True, text=True)
    if not st.stdout.strip():
        return
    subprocess.run(
        ["git", "commit", "-q", "-m", f"Build: {len(folien)} Folien"],
        cwd=WURZEL, capture_output=True)
    print("  Stand versioniert")


# --------------------------------------------------------------------------

def main():
    schnell = "--schnell" in sys.argv
    nur = None
    if "--nur" in sys.argv:
        nur = {int(x) for x in sys.argv[sys.argv.index("--nur") + 1].split(",")}

    if not QUELLE.exists():
        sys.exit(f"Fehlt: {QUELLE}")

    try:
        folien = parse(QUELLE.read_text(encoding="utf-8"))
    except SyntaxError as ex:
        sys.exit(f"\nFEHLER in folien.md\n\n{ex}\n")

    if not folien:
        sys.exit("folien.md enthaelt keine Folie (Zeile mit '## ' beginnen).")

    AUS.mkdir(exist_ok=True)
    print(f"\nKLARTEXT — {len(folien)} Folien\n")

    gesamt_haupt = len([x for x in folien if not x.get("backup")])
    lauf = 0
    for i, f in enumerate(folien, 1):
        if not f.get("backup"):
            lauf += 1
        try:
            (AUS / f"folie-{i:02d}.html").write_text(
                rendere(f, lauf, gesamt_haupt), encoding="utf-8")
        except ValueError as ex:
            sys.exit(f"\nFEHLER: {ex}\n")

    haupt = [f for f in folien if not f.get("backup")]
    n_backup = len(folien) - len(haupt)
    if len(haupt) > 15:
        print(f"  ! {len(haupt)} Folien im Hauptteil, die Aufgabe erlaubt hoechstens 15\n")
    elif n_backup:
        print(f"  {len(haupt)} Folien plus {n_backup} Backup\n")

    if schnell:
        print("  HTML geschrieben (--schnell, keine Bilder)")
        zeit(folien)
        return

    if not Path(CHROME).exists():
        sys.exit(f"Chrome nicht gefunden: {CHROME}")

    schiesse_pngs(folien, nur)
    p = baue_pptx(folien)
    d = baue_pdf_links(folien)
    print(f"\n  {p.name}")
    print(f"  {d.name}")
    zeit(folien)
    commit(folien)
    print()


if __name__ == "__main__":
    main()
